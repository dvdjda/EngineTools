"""One-slide PPTX results report for a LiBr chiller sizing case."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from .charts import make_chart, design_rows, result_rows

NAVY = RGBColor(0x2E, 0x4E, 0x7E)
TEAL = RGBColor(0x2B, 0xB6, 0xA3)
LIGHT = RGBColor(0xEA, 0xF0, 0xF8)
GREY = RGBColor(0x59, 0x59, 0x59)
INK = RGBColor(0x22, 0x30, 0x3F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Arial"


def _kv_card(slide, x, y, w, h, title, rows):
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid(); card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(0xC9, 0xD2, 0xE0); card.line.width = Pt(1)
    card.shadow.inherit = False
    tb = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.1), Inches(w - 0.3), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    rr = p.add_run(); rr.text = title
    rr.font.name = FONT; rr.font.bold = True; rr.font.size = Pt(13); rr.font.color.rgb = NAVY
    # rows
    tbl = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.55), Inches(w - 0.3), Inches(h - 0.65))
    tf = tbl.text_frame; tf.word_wrap = True
    first = True
    for (name, val, unit) in rows:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(3)
        r1 = p.add_run(); r1.text = f"{name}:  "
        r1.font.name = FONT; r1.font.size = Pt(11); r1.font.color.rgb = INK
        r2 = p.add_run(); r2.text = f"{val} {unit}"
        r2.font.name = FONT; r2.font.size = Pt(11); r2.font.bold = True; r2.font.color.rgb = NAVY


def build_pptx(dp, r, path, chart_png):
    make_chart(r, chart_png)
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # title
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    rt = p.add_run(); rt.text = "Single-effect LiBr absorption chiller  -  sizing result"
    rt.font.name = FONT; rt.font.bold = True; rt.font.size = Pt(26); rt.font.color.rgb = NAVY
    sb = slide.shapes.add_textbox(Inches(0.55), Inches(1.02), Inches(12.2), Inches(0.4))
    ps = sb.text_frame.paragraphs[0]
    rs = ps.add_run(); rs.text = "Nexa Block v1  -  screening report, verify duties vs ChemCAD"
    rs.font.name = FONT; rs.font.italic = True; rs.font.size = Pt(13); rs.font.color.rgb = GREY

    # cards
    _kv_card(slide, 0.55, 1.7, 3.6, 2.4, "Design point", design_rows(dp))
    _kv_card(slide, 0.55, 4.25, 3.6, 2.7, "Key results", result_rows(r))

    # chart image
    slide.shapes.add_picture(chart_png, Inches(4.5), Inches(2.1), width=Inches(8.3))

    # footer
    fb = slide.shapes.add_textbox(Inches(0.55), Inches(7.05), Inches(12.2), Inches(0.35))
    pf = fb.text_frame.paragraphs[0]
    rf = pf.add_run()
    rf.text = "Nexa Power Investments LLC - Confidential   |   screening report, not for certification"
    rf.font.name = FONT; rf.font.size = Pt(9); rf.font.italic = True; rf.font.color.rgb = GREY

    prs.save(path)
    return path
