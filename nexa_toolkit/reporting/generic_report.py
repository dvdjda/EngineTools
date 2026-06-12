"""
Generic reporting. Given an engine (contract) + input values + result, produce
chart, Excel, PDF and PPTX. No system-specific code here - it reads the contract,
so every registered system gets the same reporting, and upgrades land everywhere.
"""
import os
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

from .charts import NAVY, TEAL, GRID, INK

# ---- styling shared constants ----
HEX_NAVY = "2E4E7E"
HEX_LIGHT = "EAF0F8"
HEX_GREY = "595959"
BASIS_COLOR = {"verified": "2E7D4E", "screening": "B26A00", "input": "595959", "unverified": "C0392B"}


def _design_rows(engine, values):
    return [(i.label, f"{values[i.key]:g}", i.unit) for i in engine.inputs]


def _convergence_of(result):
    """Return the SolvedSystem's ConvergenceStatus if result carries one, else None.
    v1 engines (no 'solved' key) → None → existing report layout, unchanged."""
    solved = result.get("solved") if isinstance(result, dict) else None
    return getattr(solved, "convergence", None) if solved is not None else None


def _feasibility_of(result):
    """Return result['feasibility'] (FeasibilityStatus) if the engine surfaces it.
    Strictly independent of convergence — different failure mode."""
    return result.get("feasibility") if isinstance(result, dict) else None


def _any_failed(result) -> bool:
    """Either solver or feasibility check failed → KPIs untrustworthy."""
    conv = _convergence_of(result)
    feas = _feasibility_of(result)
    return ((conv is not None and not conv.converged) or
            (feas is not None and not feas.feasible))


def _result_rows(engine, result):
    """Result rows. Override basis to 'unverified' (red) when either the
    solver loops failed OR the power balance is infeasible — both are
    independent reasons to distrust the KPIs."""
    not_ok = _any_failed(result)
    return [(o.label, o.text(), o.unit,
             "unverified" if not_ok else o.basis)
            for o in engine.outputs(result)]


def _convergence_text(result):
    """One-line summary + ok-flag. (None, True) if the result has no convergence info."""
    from nexablock.core.convergence import convergence_summary
    conv = _convergence_of(result)
    if conv is None:
        return None, True
    return convergence_summary(conv)


# ---------- shared AI text parser ----------
def _parse_ai_sections(ai_text):
    """Split '1. HEADING body' LLM text into [(heading, body), ...] pairs."""
    import re
    parts = re.split(r'(\d\.\s+[A-Z][A-Z &]+)', ai_text or "")
    out = []
    i = 0
    while i < len(parts):
        chunk = parts[i].strip()
        if not chunk:
            i += 1; continue
        if re.match(r'^\d\.\s+[A-Z]', chunk):
            body = parts[i + 1].strip() if i + 1 < len(parts) else ""
            out.append((chunk, body))
            i += 2
        else:
            out.append(("", chunk))
            i += 1
    return out


# ---------- chart ----------
def write_study_sheet(wb, study, sheet_name: str = "Study", image_anchor: str = "A3",
                      table_start_row: int = 24):
    """Add a Study sheet to `wb`: navy banner + embedded chart + table.

    Shared by build_excel(..., study=...) and the standalone study_to_xlsx
    in nexa_toolkit.reporting.study_export — one source of truth for what
    a "study sheet" looks like.

    Returns the worksheet.
    """
    from openpyxl.styles import Font, PatternFill, Alignment
    ws = wb.create_sheet(sheet_name)
    ws["A1"] = study.get("label") or "Study"
    ws["A1"].font = Font(name="Arial", bold=True, size=12, color="FFFFFF")
    ws["A1"].fill = PatternFill("solid", fgColor=HEX_NAVY)
    ws["A1"].alignment = Alignment(vertical="center", indent=1)
    ws.column_dimensions["A"].width = 22
    for col in "BCDEFGHIJKLMN":
        ws.column_dimensions[col].width = 16

    # Chart (PNG written next to the workbook path)
    try:
        # The png is written to a temp file with a stable name so concurrent
        # exports don't collide.
        import tempfile
        study_png = os.path.join(tempfile.mkdtemp(), f"{sheet_name}.png")
        _render_study_chart(study, study_png)
        from openpyxl.drawing.image import Image as XImage
        ws.add_image(XImage(study_png), image_anchor)
    except Exception as _e:
        ws[image_anchor] = f"Study chart omitted: {_e}"

    # Table from result.as_dataframe(). Pandas may be absent.
    ws.cell(table_start_row, 1, "Study data").font = Font(bold=True)
    try:
        df = study["result"].as_dataframe()
        if getattr(df.index, "name", None) or df.index.dtype == object:
            df = df.reset_index().rename(columns={"index": "scenario"})
        headers = [str(c) for c in df.columns]
        for j, h in enumerate(headers):
            ws.cell(table_start_row + 1, j + 1, h).font = Font(bold=True)
        for i, row in enumerate(df.values.tolist()):
            for j, val in enumerate(row):
                if isinstance(val, float) and (val != val):     # NaN
                    ws.cell(table_start_row + 2 + i, j + 1, None)
                else:
                    ws.cell(table_start_row + 2 + i, j + 1, val)
    except ImportError:
        ws.cell(table_start_row + 1, 1, "table omitted: pandas not installed")
    except Exception as _e:
        ws.cell(table_start_row + 1, 1, f"table omitted: {_e}")
    return ws


def _render_study_chart(study, path):
    """Dispatch to the matching nexablock.studies chart helper.

    Accepts:  study = {"kind": "sensitivity"|"sweep"|"scenarios",
                       "result": <result>,
                       "chart_kwargs": dict,
                       "label": str}
    Returns:  path to the PNG written. Re-uses the chart helpers shipped
              in nexablock.studies.charts — no second chart-generation path.
    """
    from nexablock.studies import tornado_chart, sweep_chart, scenarios_chart
    fns = {"sensitivity": tornado_chart,
           "sweep":       sweep_chart,
           "scenarios":   scenarios_chart}
    fn = fns.get(study["kind"])
    if fn is None:
        raise ValueError(f"Unknown study kind: {study['kind']!r}")
    fn(study["result"], path, **study.get("chart_kwargs", {}))
    return path


def build_chart(engine, result, path):
    # If the engine emits SVG but the caller wants a raster (PDF/PPTX),
    # let the engine write SVG to a sibling file and rasterize via cairosvg.
    fmt = getattr(engine, "chart_format", "png")
    target_ext = os.path.splitext(path)[1].lower()
    if fmt == "svg" and target_ext != ".svg":
        import cairosvg
        svg_path = path + ".tmp.svg"
        engine.chart(result, svg_path)
        cairosvg.svg2png(url=svg_path, write_to=path,
                         output_width=1600, background_color="white")
        return path
    sig = engine.chart(result, path)
    if sig:
        return sig
    # generic fallback: bar of numeric highlight outputs
    hl = engine.highlights(result)
    fig, ax = plt.subplots(figsize=(9.2, 3.5))
    fig.subplots_adjust(left=0.1, right=0.97, top=0.85, bottom=0.18)
    labels = [o.label for o in hl]
    vals = [o.value for o in hl]
    ax.bar(labels, vals, color=NAVY, width=0.5, zorder=3)
    ax.set_title(engine.name)
    ax.yaxis.grid(True, color=GRID); ax.set_axisbelow(True)
    ax.spines[["top", "right"]].set_visible(False)
    for i, v in enumerate(vals):
        ax.text(i, v, f"{v:g}", ha="center", va="bottom", fontsize=9, color=INK)
    fig.savefig(path, dpi=150, facecolor="white")
    plt.close(fig)
    return path


# ---------- Excel ----------
def build_excel(engine, values, result, path, ai_text=None, study=None):
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.chart import BarChart, Reference
    thin = Side(style="thin", color="C9D2E0")
    box = Border(left=thin, right=thin, top=thin, bottom=thin)
    wb = Workbook(); ws = wb.active; ws.title = "Sizing"
    ws.sheet_view.showGridLines = False
    for col, w in (("A", 30), ("B", 14), ("C", 12), ("D", 12)):
        ws.column_dimensions[col].width = w

    ws["A1"] = engine.name
    ws["A1"].font = Font(name="Arial", bold=True, size=14, color=HEX_NAVY)
    ws["A2"] = "Screening report  -  verify vs ChemCAD where noted"
    ws["A2"].font = Font(name="Arial", italic=True, size=10, color=HEX_GREY)

    def band(cell, text, span):
        ws[cell] = text
        ws[cell].font = Font(name="Arial", bold=True, size=12, color="FFFFFF")
        ws[cell].fill = PatternFill("solid", fgColor=HEX_NAVY)
        ws[cell].alignment = Alignment(vertical="center", indent=1)
        ws.merge_cells(f"{cell}:{span}")

    def header(row, cols):
        for col, t in zip("ABCD", cols):
            c = ws[f"{col}{row}"]; c.value = t
            c.font = Font(name="Arial", bold=True, color=HEX_NAVY)
            c.fill = PatternFill("solid", fgColor=HEX_LIGHT); c.border = box
            c.alignment = Alignment(horizontal="left" if col == "A" else "center")

    band("A4", "Design point (inputs)", "C4")
    header(5, ("Quantity", "Value", "Unit"))
    row = 6
    for name, val, unit in _design_rows(engine, values):
        ws[f"A{row}"] = name
        ws[f"B{row}"] = val; ws[f"B{row}"].font = Font(name="Arial", color="0000FF")
        ws[f"C{row}"] = unit
        for col in "ABC":
            ws[f"{col}{row}"].border = box
            if col != "A":
                ws[f"{col}{row}"].alignment = Alignment(horizontal="center")
        row += 1

    # Convergence status row(s) — only present when result carries a SolvedSystem.
    conv_text, conv_ok = _convergence_text(result)
    if conv_text is not None:
        crow = row + 2
        band(f"A{crow}", "Convergence", f"D{crow}")
        ws.cell(crow + 1, 1, conv_text).font = Font(
            name="Arial", bold=True, size=11,
            color=("2E7D4E" if conv_ok else "C0392B"))
        ws.merge_cells(f"A{crow + 1}:D{crow + 1}")
        if not conv_ok:
            ws.cell(crow + 2, 1,
                    "⚠ NOT CONVERGED — KPIs below may be unreliable.").font = Font(
                name="Arial", bold=True, size=10, color="C0392B")
            ws.merge_cells(f"A{crow + 2}:D{crow + 2}")
            row = crow + 3
        else:
            row = crow + 2

    # Resource balance bands — one per ResourceBalance.
    feas = _feasibility_of(result)
    if feas is not None and getattr(feas, "balances", None):
        for b in feas.balances:
            frow = row + 2
            band(f"A{frow}", f"{b.resource} balance", f"D{frow}")
            ws.cell(frow + 1, 1, f"Assumption: {b.assumption}").font = Font(
                name="Arial", italic=True, size=10, color=HEX_GREY)
            ws.merge_cells(f"A{frow + 1}:D{frow + 1}")
            for i, (lab, val) in enumerate([
                ("Supply",    b.supply),
                ("Demand",    b.demand),
                ("Balance",   b.balance),
                ("Shortfall", b.shortfall),
            ], start=frow + 2):
                ws.cell(i, 1, lab).font = Font(name="Arial", bold=True)
                ws.cell(i, 2, round(val, 1))
                ws.cell(i, 3, b.unit)
                if lab == "Balance":
                    ws.cell(i, 2).font = Font(
                        name="Arial", bold=True,
                        color=("2E7D4E" if b.feasible else "C0392B"))
            bdrow = frow + 6
            ws.cell(bdrow, 1, "Breakdown").font = Font(
                name="Arial", bold=True, color=HEX_GREY)
            for i, (k, v) in enumerate(b.breakdown.items(), start=bdrow + 1):
                ws.cell(i, 1, k)
                ws.cell(i, 2, "not modelled" if v is None else round(v, 1))
                ws.cell(i, 3, "" if v is None else b.unit)
            next_row = bdrow + 1 + len(b.breakdown)
            if not b.feasible:
                ws.cell(next_row + 1, 1,
                        f"⚠ {b.resource.upper()} DEFICIT — demand {b.demand:,.0f} "
                        f"{b.unit} > supply {b.supply:,.0f} {b.unit}, "
                        f"shortfall {b.shortfall:,.0f} {b.unit}.").font = Font(
                    name="Arial", bold=True, size=11, color="C0392B")
                ws.merge_cells(f"A{next_row + 1}:D{next_row + 1}")
                row = next_row + 2
            else:
                row = next_row

    start = row + 2
    band(f"A{start}", "Results", f"D{start}")
    header(start + 1, ("Quantity", "Value", "Unit", "Basis"))
    row = start + 2
    kw_rows = []
    for label, value, unit, basis in _result_rows(engine, result):
        ws[f"A{row}"] = label
        ws[f"B{row}"] = value
        ws[f"C{row}"] = unit
        ws[f"D{row}"] = basis
        ws[f"D{row}"].font = Font(name="Arial", size=10, color=BASIS_COLOR.get(basis, "595959"))
        if unit == "kW":
            kw_rows.append(row)
        for col in "ABCD":
            ws[f"{col}{row}"].border = box
            if col in "BCD":
                ws[f"{col}{row}"].alignment = Alignment(horizontal="center")
        row += 1

    if kw_rows:
        chart = BarChart(); chart.type = "col"; chart.title = "Heat / duties (kW)"
        chart.legend = None; chart.height = 7; chart.width = 12
        data = Reference(ws, min_col=2, min_row=min(kw_rows), max_row=max(kw_rows))
        cats = Reference(ws, min_col=1, min_row=min(kw_rows), max_row=max(kw_rows))
        chart.add_data(data); chart.set_categories(cats)
        ws.add_chart(chart, "F4")

    # --- AI Analysis section ---
    if ai_text:
        row += 1
        band(f"A{row}", "AI Engineering Analysis", f"D{row}")
        row += 1
        sub_cell = ws[f"A{row}"]
        sub_cell.value = "Multidisciplinary expert assessment  -  screening only"
        sub_cell.font = Font(name="Arial", italic=True, size=9, color=HEX_GREY)
        ws.merge_cells(f"A{row}:D{row}")
        row += 1
        for heading, body in _parse_ai_sections(ai_text):
            if heading:
                ws[f"A{row}"] = heading
                ws[f"A{row}"].font = Font(name="Arial", bold=True, size=10, color=HEX_NAVY)
                ws[f"A{row}"].alignment = Alignment(indent=1)
                ws.merge_cells(f"A{row}:D{row}")
                row += 1
            if body:
                ws[f"A{row}"] = body
                ws[f"A{row}"].font = Font(name="Arial", size=9, color="222222")
                ws[f"A{row}"].alignment = Alignment(wrap_text=True, vertical="top", indent=2)
                ws[f"A{row}"].border = box
                ws.merge_cells(f"A{row}:D{row}")
                ws.row_dimensions[row].height = max(45, min(200, len(body) // 3))
                row += 1

    for r_ in ws.iter_rows():
        for c in r_:
            if c.value is not None and (not c.font or c.font.name != "Arial"):
                f = c.font
                c.font = Font(name="Arial", bold=f.bold, italic=f.italic, size=f.size, color=f.color)
    # Optional Study sheet — single source of truth in write_study_sheet.
    if study is not None:
        write_study_sheet(wb, study)

    wb.save(path)
    return path


# ---------- PDF ----------
def build_pdf(engine, values, result, path, chart_png, ai_text=None, study=None):
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import mm
    from reportlab.lib import colors
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer, Table,
                                    TableStyle, Image)
    from reportlab.lib.enums import TA_LEFT
    navy = colors.HexColor("#2E4E7E"); grey = colors.HexColor("#595959")
    ink = colors.HexColor("#22303F"); light = colors.HexColor("#EAF0F8")
    bcol = {"verified": colors.HexColor("#2E7D4E"),
            "screening": colors.HexColor("#B26A00"),
            "input": colors.HexColor("#595959"), "unverified": colors.HexColor("#C0392B")}
    build_chart(engine, result, chart_png)
    ss = getSampleStyleSheet()
    h1 = ParagraphStyle("h1", parent=ss["Title"], fontName="Helvetica-Bold", fontSize=17,
                        textColor=navy, alignment=TA_LEFT, spaceAfter=12, leading=22)
    sub = ParagraphStyle("sub", fontName="Helvetica-Oblique", fontSize=9.5, textColor=grey, spaceAfter=10)
    sec = ParagraphStyle("sec", fontName="Helvetica-Bold", fontSize=11, textColor=navy,
                         spaceBefore=8, spaceAfter=4)
    body = ParagraphStyle("body", fontName="Helvetica", fontSize=8.5, textColor=ink, leading=11.5)
    basis_st = {b: ParagraphStyle(f"b{b}", fontName="Helvetica", fontSize=8, textColor=c)
                for b, c in bcol.items()}

    doc = SimpleDocTemplate(path, pagesize=A4, leftMargin=18 * mm, rightMargin=18 * mm,
                            topMargin=16 * mm, bottomMargin=16 * mm)
    W = doc.width
    story = [Paragraph(engine.name, h1),
             Paragraph("Nexa Block v1  &nbsp;|&nbsp;  screening report", sub)]

    # design point table
    story.append(Paragraph("Design point", sec))
    d = Table([[n, v, u] for (n, v, u) in _design_rows(engine, values)],
              colWidths=[W * 0.55, W * 0.18, W * 0.27])
    d.setStyle(TableStyle([
        ("FONT", (0, 0), (-1, -1), "Helvetica", 9), ("TEXTCOLOR", (0, 0), (-1, -1), ink),
        ("ALIGN", (1, 0), (-1, -1), "CENTER"),
        ("ROWBACKGROUNDS", (0, 0), (-1, -1), [colors.white, light]),
        ("LINEBELOW", (0, 0), (-1, -1), 0.4, colors.HexColor("#C9D2E0")),
        ("TOPPADDING", (0, 0), (-1, -1), 4), ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
        ("LEFTPADDING", (0, 0), (0, -1), 8)]))
    story.append(d)

    # Convergence status (only if the result carries a SolvedSystem).
    conv_text, conv_ok = _convergence_text(result)
    if conv_text is not None:
        story.append(Spacer(1, 8))
        story.append(Paragraph("Convergence", sec))
        good_st = ParagraphStyle("conv_ok",  fontName="Helvetica-Bold", fontSize=9,
                                  textColor=colors.HexColor("#2E7D4E"))
        bad_st  = ParagraphStyle("conv_bad", fontName="Helvetica-Bold", fontSize=9.5,
                                  textColor=colors.HexColor("#C0392B"))
        story.append(Paragraph(conv_text, good_st if conv_ok else bad_st))
        if not conv_ok:
            warn_st = ParagraphStyle("conv_warn", fontName="Helvetica-Bold",
                                      fontSize=10,
                                      textColor=colors.HexColor("#C0392B"),
                                      backColor=colors.HexColor("#FBE9E7"),
                                      borderPadding=6, borderRadius=4,
                                      spaceBefore=4, spaceAfter=4)
            story.append(Paragraph(
                "⚠ KPIs below may be unreliable — solve did not converge.",
                warn_st))

    # Resource balance sections — one per ResourceBalance, strictly separate
    # from convergence.
    feas = _feasibility_of(result)
    if feas is not None and getattr(feas, "balances", None):
        green = colors.HexColor("#2E7D4E"); red_c = colors.HexColor("#C0392B")
        assum_st = ParagraphStyle("feas_assum", fontName="Helvetica-Oblique",
                                   fontSize=8.5, textColor=grey, spaceAfter=4)
        for b in feas.balances:
            story.append(Spacer(1, 8))
            story.append(Paragraph(f"{b.resource} balance", sec))
            story.append(Paragraph(f"Assumption: {b.assumption}", assum_st))
            bal_color = green if b.feasible else red_c
            bal_st = ParagraphStyle(f"feas_bal_{b.resource}", fontName="Helvetica-Bold",
                                     fontSize=10, textColor=bal_color)
            story.append(Paragraph(
                f"Supply {b.supply:,.0f} {b.unit} · "
                f"Demand {b.demand:,.0f} {b.unit} · "
                f"Balance {b.balance:+,.0f} {b.unit}", bal_st))
            bd_rows = [[k, "not modelled" if v is None else f"{v:+,.1f} {b.unit}"]
                       for k, v in b.breakdown.items()]
            if bd_rows:
                bd = Table(bd_rows, colWidths=[W * 0.55, W * 0.35])
                bd.setStyle(TableStyle([
                    ("FONT", (0, 0), (-1, -1), "Helvetica", 8.5),
                    ("TEXTCOLOR", (0, 0), (-1, -1), ink),
                    ("ROWBACKGROUNDS", (0, 0), (-1, -1), [colors.white, light]),
                    ("ALIGN", (1, 0), (-1, -1), "RIGHT"),
                    ("LEFTPADDING", (0, 0), (0, -1), 8),
                    ("TOPPADDING", (0, 0), (-1, -1), 2),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 2)]))
                story.append(bd)
            if not b.feasible:
                warn_st = ParagraphStyle(f"feas_warn_{b.resource}",
                                          fontName="Helvetica-Bold", fontSize=10.5,
                                          textColor=red_c,
                                          backColor=colors.HexColor("#FBE9E7"),
                                          borderPadding=6, borderRadius=4,
                                          spaceBefore=6, spaceAfter=6)
                resource_upper = b.resource.upper()
                story.append(Paragraph(
                    f"⚠ {resource_upper} DEFICIT — demand {b.demand:,.0f} "
                    f"{b.unit} &gt; supply {b.supply:,.0f} {b.unit}, "
                    f"shortfall {b.shortfall:,.0f} {b.unit}.", warn_st))

    # results table with basis column
    story.append(Paragraph("Results", sec))
    rdata = [[Paragraph("<b>Quantity</b>", body), Paragraph("<b>Value</b>", body),
              Paragraph("<b>Unit</b>", body), Paragraph("<b>Basis</b>", body)]]
    for label, value, unit, basis in _result_rows(engine, result):
        rdata.append([Paragraph(label, body), Paragraph(value, body),
                      Paragraph(unit, body), Paragraph(basis, basis_st.get(basis, body))])
    rt = Table(rdata, colWidths=[W * 0.46, W * 0.18, W * 0.18, W * 0.18])
    rt.setStyle(TableStyle([
        ("ALIGN", (1, 0), (-1, -1), "CENTER"),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, light]),
        ("LINEBELOW", (0, 0), (-1, -1), 0.4, colors.HexColor("#C9D2E0")),
        ("TOPPADDING", (0, 0), (-1, -1), 3), ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
        ("LEFTPADDING", (0, 0), (0, -1), 6)]))
    story.append(rt)

    story.append(Paragraph("Chart", sec))
    story.append(Image(chart_png, width=W, height=W * 3.5 / 9.2))
    story.append(Spacer(1, 6))

    # Optional study chart appended after the flowsheet chart.
    if study is not None:
        try:
            study_png = chart_png + ".study.png"
            _render_study_chart(study, study_png)
            story.append(Spacer(1, 12))
            story.append(Paragraph(study.get("label") or "Study", sec))
            story.append(Image(study_png, width=W, height=W * 3.5 / 9.2))
            story.append(Spacer(1, 6))
        except Exception as _e:
            story.append(Paragraph(f"<i>Study chart omitted: {_e}</i>", body))

    if engine.notes:
        story.append(Paragraph("Method: " + engine.notes, body))

    # --- AI Analysis section ---
    if ai_text:
        story.append(Spacer(1, 10))
        story.append(Paragraph("AI Engineering Analysis", sec))
        story.append(Paragraph(
            "Multidisciplinary expert assessment  \u2014  screening only, not for certification.", sub))
        for heading, body_txt in _parse_ai_sections(ai_text):
            if heading:
                story.append(Paragraph(f"<b>{heading}</b>", body))
            if body_txt:
                story.append(Paragraph(body_txt, body))
                story.append(Spacer(1, 5))

    def footer(canvas, d_):
        canvas.saveState()
        canvas.setStrokeColor(navy); canvas.setLineWidth(2)
        canvas.line(18 * mm, A4[1] - 12 * mm, A4[0] - 18 * mm, A4[1] - 12 * mm)
        canvas.setFont("Helvetica", 7.5); canvas.setFillColor(grey)
        canvas.drawString(18 * mm, 10 * mm, "Nexa Power Investments LLC - Confidential")
        canvas.drawRightString(A4[0] - 18 * mm, 10 * mm, "Screening report - not for certification")
        canvas.restoreState()
    doc.build(story, onFirstPage=footer, onLaterPages=footer)
    return path


# ---------- PPTX ----------
def _pptx_textbox(slide, x, y, w, h, runs):
    """Add a textbox to a pptx slide. runs = list of lines; line = [(text,size,bold,RGBColor)]."""
    from pptx.util import Inches, Pt
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; first = True
    for line in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.space_after = Pt(3)
        for (t, sz, b, col) in line:
            r = p.add_run(); r.text = t
            r.font.name = "Arial"; r.font.size = Pt(sz); r.font.bold = b; r.font.color.rgb = col
    return tb


def build_pptx(engine, values, result, path, chart_png, ai_text=None):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    navy  = RGBColor(0x2E, 0x4E, 0x7E); grey  = RGBColor(0x59, 0x59, 0x59)
    ink   = RGBColor(0x22, 0x30, 0x3F); white = RGBColor(0xFF, 0xFF, 0xFF)
    build_chart(engine, result, chart_png)
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)

    # ---- Slide 1: results + chart ----
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    def txt(x, y, w, h, runs):
        return _pptx_textbox(slide, x, y, w, h, runs)

    def card(x, y, w, h, title, rows):
        c = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        c.fill.solid(); c.fill.fore_color.rgb = white
        c.line.color.rgb = RGBColor(0xC9, 0xD2, 0xE0); c.line.width = Pt(1); c.shadow.inherit = False
        txt(x + 0.15, y + 0.1, w - 0.3, 0.4, [[(title, 13, True, navy)]])
        lines = [[(f"{n}:  ", 11, False, ink), (f"{v} {u}", 11, True, navy)] for (n, v, u) in rows]
        txt(x + 0.15, y + 0.55, w - 0.3, h - 0.6, lines)

    txt(0.55, 0.35, 12.2, 0.70, [[(engine.name, 24, True, navy)]])
    txt(0.55, 1.05, 12.2, 0.40, [[("Nexa Block v1  -  screening report", 13, False, grey)]])
    card(0.55, 1.70, 3.7, 2.5, "Design point",
         [(n, v, u) for (n, v, u) in _design_rows(engine, values)])
    hl = [(o.label, o.text(), o.unit) for o in engine.highlights(result)]
    card(0.55, 4.35, 3.7, 2.4, "Key results", hl)
    slide.shapes.add_picture(chart_png, Inches(4.55), Inches(2.3), width=Inches(8.25))
    txt(0.55, 7.05, 12.2, 0.35,
        [[("Nexa Power Investments LLC - Confidential   |   screening report, not for certification",
           9, False, grey)]])

    # ---- Slide 2: AI Analysis (only when AI has been run) ----
    if ai_text:
        slide2 = prs.slides.add_slide(prs.slide_layouts[6])

        def txt2(x, y, w, h, runs):
            return _pptx_textbox(slide2, x, y, w, h, runs)

        txt2(0.45, 0.18, 12.4, 0.65,
             [[(f"AI Engineering Analysis  -  {engine.name}", 20, True, navy)]])
        txt2(0.45, 0.85, 12.4, 0.38,
             [[("Multidisciplinary expert assessment  -  screening only, not for certification",
                11, False, grey)]])

        sections = _parse_ai_sections(ai_text)
        half = max(1, (len(sections) + 1) // 2)
        left_secs, right_secs = sections[:half], sections[half:]

        def _sec_runs(secs):
            runs = []
            for heading, body in secs:
                if heading:
                    runs.append([(heading, 10, True, navy)])
                if body:
                    runs.append([(body, 9, False, ink)])
                runs.append([("", 5, False, ink)])  # spacer line
            return runs

        txt2(0.45, 1.35, 5.90, 5.85, _sec_runs(left_secs))
        if right_secs:
            txt2(7.00, 1.35, 5.90, 5.85, _sec_runs(right_secs))

        txt2(0.45, 7.10, 12.4, 0.32,
             [[("Nexa Power Investments LLC - Confidential   |   AI analysis - screening only",
                9, False, grey)]])

    prs.save(path)
    return path


def build_all(engine, values, outdir, slug=None):
    os.makedirs(outdir, exist_ok=True)
    slug = slug or engine.key
    result = engine.solve(values)
    build_chart(engine, result, f"{outdir}/{slug}_chart.png")
    build_excel(engine, values, result, f"{outdir}/{slug}.xlsx")
    build_pdf(engine, values, result, f"{outdir}/{slug}.pdf", f"{outdir}/_{slug}_pdf.png")
    build_pptx(engine, values, result, f"{outdir}/{slug}.pptx", f"{outdir}/_{slug}_pptx.png")
    return result
